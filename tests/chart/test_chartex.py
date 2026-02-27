# pyright: reportPrivateUsage=false

"""Test suite for `pptx.chart.chartex` module and `WaterfallChartData`."""

from __future__ import annotations

import zipfile
from io import BytesIO

import pytest

from pptx.chart.data import WaterfallChartData
from pptx.util import Inches


class DescribeWaterfallChartData:
    """Unit-test suite for `pptx.chart.data.WaterfallChartData`."""

    def it_can_set_categories(self):
        chart_data = WaterfallChartData()
        chart_data.categories = ["Q1", "Q2", "Q3", "Total"]
        assert chart_data.categories == ["Q1", "Q2", "Q3", "Total"]

    def it_can_add_a_series(self):
        chart_data = WaterfallChartData()
        chart_data.categories = ["Q1", "Q2", "Total"]
        chart_data.add_series("Revenue", [100, 50, 150], subtotals=[2])

        assert chart_data.series_name == "Revenue"
        assert chart_data.series_values == [100, 50, 150]
        assert chart_data.subtotals == [2]

    def it_provides_excel_refs(self):
        chart_data = WaterfallChartData()
        chart_data.categories = ["A", "B", "C"]
        chart_data.add_series("Sales", [1, 2, 3])

        assert chart_data.categories_ref == "Sheet1!$A$2:$A$4"
        assert chart_data.values_ref == "Sheet1!$B$2:$B$4"
        assert chart_data.series_name_ref == "Sheet1!$B$1"

    def it_can_generate_an_xlsx_blob(self):
        chart_data = WaterfallChartData()
        chart_data.categories = ["Q1", "Q2"]
        chart_data.add_series("Revenue", [100, 200])

        xlsx_blob = chart_data.xlsx_blob

        assert isinstance(xlsx_blob, bytes)
        assert len(xlsx_blob) > 0
        # verify it's a valid zip (xlsx is a zip archive)
        zf = zipfile.ZipFile(BytesIO(xlsx_blob))
        assert "xl/worksheets/sheet1.xml" in zf.namelist()
        zf.close()

    def it_raises_on_mismatched_categories_and_values(self):
        chart_data = WaterfallChartData()
        chart_data.categories = ["Q1", "Q2", "Q3"]
        chart_data.add_series("Revenue", [100, 200])

        with pytest.raises(ValueError, match="categories length.*must equal.*series values"):
            chart_data.xlsx_blob


class DescribeChartEx:
    """Unit-test suite for `pptx.chart.chartex.ChartEx`."""

    def it_applies_subtotals_from_chart_data(self):
        """Subtotals set via chart_data.add_series flow through to chart XML."""
        from pptx import Presentation
        from pptx.oxml.ns import qn

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        chart_data = WaterfallChartData()
        chart_data.categories = ["Q1", "Q2", "Q3", "Total"]
        chart_data.add_series("Revenue", [100, 50, -30, 120], subtotals=[3])

        graphic_frame = slide.shapes.add_chartex(
            chart_data, Inches(1), Inches(1), Inches(6), Inches(4),
        )

        chartex = graphic_frame.chartex
        # verify the subtotal idx=3 appears in the XML
        series_el = chartex._element.chart.plotArea.plotAreaRegion.findall(
            qn("cx:series")
        )[0]
        subtotals = series_el.findall(f".//{qn('cx:subtotals')}/{qn('cx:idx')}")
        assert [int(el.get("val")) for el in subtotals] == [3]

    def it_can_replace_data(self):
        """End-to-end: build via add_chartex → replace_data → verify."""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        chart_data = WaterfallChartData()
        chart_data.categories = ["Q1", "Q2", "Q3", "Total"]
        chart_data.add_series("Revenue", [100, 50, -30, 120], subtotals=[3])

        graphic_frame = slide.shapes.add_chartex(
            chart_data,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(4),
        )

        chartex = graphic_frame.chartex
        assert chartex.chart_type == "waterfall"

        series_list = chartex.series
        assert len(series_list) == 1
        assert series_list[0].name == "Revenue"
        assert series_list[0].values == [100.0, 50.0, -30.0, 120.0]

        # --- now replace the data and verify ---
        new_data = WaterfallChartData()
        new_data.categories = ["Jan", "Feb", "Mar", "Apr", "Total"]
        new_data.add_series("Profit", [200, -50, 100, 75, 325], subtotals=[4])

        chartex.replace_data(new_data)

        # re-fetch series from the live XML
        chartex2 = graphic_frame.chartex
        series_list2 = chartex2.series
        assert len(series_list2) == 1
        assert series_list2[0].name == "Profit"
        assert series_list2[0].values == [200.0, -50.0, 100.0, 75.0, 325.0]

    def it_removes_stale_dataPt_on_replace_data(self):
        """replace_data with fewer categories removes out-of-range dataPt elements."""
        from pptx import Presentation
        from pptx.oxml.ns import qn
        from pptx.oxml.xmlchemy import OxmlElement

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        chart_data = WaterfallChartData()
        chart_data.categories = ["A", "B", "C", "D", "E", "F", "G", "H"]
        chart_data.add_series("S", [1, 2, 3, 4, 5, 6, 7, 8], subtotals=[7])

        graphic_frame = slide.shapes.add_chartex(
            chart_data, Inches(1), Inches(1), Inches(6), Inches(4),
        )
        chartex = graphic_frame.chartex

        # --- manually add dataPt elements to the series (simulating template) ---
        series_el = chartex._element.chart.plotArea.plotAreaRegion.series_lst[0]
        for i in range(8):
            dataPt = OxmlElement("cx:dataPt")
            dataPt.set("idx", str(i))
            series_el.append(dataPt)

        assert len(series_el.dataPt_lst) == 8

        # --- replace with fewer categories ---
        new_data = WaterfallChartData()
        new_data.categories = ["X", "Y", "Z", "Total"]
        new_data.add_series("S2", [10, 20, 30, 60], subtotals=[3])

        chartex.replace_data(new_data)

        remaining = series_el.dataPt_lst
        remaining_indices = [int(dp.get("idx")) for dp in remaining]
        # only indices 0-3 should remain (4 categories)
        assert all(idx < 4 for idx in remaining_indices)
        # indices 4-7 should have been removed
        assert not any(idx >= 4 for idx in remaining_indices)
